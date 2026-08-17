#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""evo-kernel 双 Harness 系统 —— 架构与工程开发能力 PPT 生成器
暗色科技风，16:9，python-pptx。输出 .pptx（Keynote/PowerPoint 可开）。"""
import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 调色板（对齐 system-overview.html 暗色主题） ──────────────────────────
BG      = RGBColor(0x0B, 0x0E, 0x14)   # 页面底
PANEL   = RGBColor(0x14, 0x18, 0x24)   # 面板
CARD    = RGBColor(0x18, 0x1E, 0x2D)   # 卡片
EDGE    = RGBColor(0x23, 0x2B, 0x3B)   # 卡片描边
CHEV    = RGBColor(0x20, 0x2B, 0x42)   # chevron 流程块填充
INK     = RGBColor(0xE8, 0xED, 0xF5)   # 主文字
INK2    = RGBColor(0x9A, 0xA7, 0xB8)   # 次文字
INK3    = RGBColor(0x6B, 0x76, 0x88)   # 弱文字
ACCENT  = RGBColor(0x4D, 0x9D, 0xFF)   # 蓝（主强调）
ACC_D   = RGBColor(0x2A, 0x78, 0xD6)   # 深蓝
GREEN   = RGBColor(0x3E, 0xC0, 0x3E)   # 绿（注入/生效）
AMBER   = RGBColor(0xF0, 0xB5, 0x3C)   # 琥珀（人审）
RED     = RGBColor(0xEF, 0x83, 0x83)   # 红（guard/阻断）
VIOLET  = RGBColor(0x90, 0x85, 0xE9)   # 紫（facts/语义）
CYAN    = RGBColor(0x5F, 0xC8, 0xD8)   # 青

CN   = "PingFang SC"
MONO = "Menlo"

EMU_IN = 914400

def _set_font(run, name):
    run.font.name = name                       # a:latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)

def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75, round=False, radius=0.10):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def shape(slide, kind, x, y, w, h, fill=None, line=None, lw=0.75):
    shp = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp

def tb(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    return tf

def para(tf, first=False, align=PP_ALIGN.LEFT, space_before=0, space_after=0, line=1.0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line
    return p

def run(p, text, size=14, color=INK, bold=False, font=CN, italic=False):
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color
    _set_font(r, font)
    return r

def text(slide, x, y, w, h, s, size=14, color=INK, bold=False, font=CN,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.0):
    tf = tb(slide, x, y, w, h, anchor=anchor)
    p = para(tf, first=True, align=align, line=line)
    run(p, s, size=size, color=color, bold=bold, font=font)
    return tf

# ── 页眉 / 页脚 ────────────────────────────────────────────────────────────
TOTAL = 10
def header(slide, eyebrow, title, idx):
    text(slide, 0.75, 0.40, 11.0, 0.3, eyebrow, size=10.5, color=INK3, font=MONO)
    text(slide, 0.73, 0.68, 11.8, 0.75, title, size=29, color=INK, bold=True)
    rect(slide, 0.78, 1.42, 0.52, 0.05, fill=ACCENT)
    # 右上角章节编号
    text(slide, 12.05, 0.42, 0.55, 0.3, f"{idx:02d}", size=12, color=INK3, font=MONO, align=PP_ALIGN.RIGHT)

def footer(slide, idx):
    text(slide, 0.75, 7.10, 6.0, 0.28, "evo-kernel · 经验内核", size=9.5, color=INK3, font=MONO)
    text(slide, 11.55, 7.10, 1.0, 0.28, f"{idx} / {TOTAL}", size=9.5, color=INK3, font=MONO, align=PP_ALIGN.RIGHT)

def chip(slide, x, y, w, h, s, color=ACCENT, size=10.5, fill=CARD, line=None, font=MONO, bold=False):
    c = rect(slide, x, y, w, h, fill=fill, line=line if line else color, lw=1.0, round=True, radius=0.5)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run(p, s, size=size, color=color, bold=bold, font=font)
    return c

def card(slide, x, y, w, h, edge=EDGE, fill=CARD):
    return rect(slide, x, y, w, h, fill=fill, line=edge, lw=1.0, round=True, radius=0.06)

# ═══════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(BLANK)

# ── S1 封面 ────────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
rect(s, 0, 0, 13.333, 0.16, fill=ACC_D)                       # 顶条
text(s, 1.0, 1.55, 11.3, 0.35, "EXPERIENCE GOVERNANCE KERNEL", size=13, color=ACCENT, font=MONO)
# 大标题
text(s, 0.95, 2.05, 11.4, 1.5, "evo-kernel", size=66, color=INK, bold=True, font=MONO)
rect(s, 1.0, 3.28, 0.7, 0.06, fill=GREEN)
text(s, 1.0, 3.55, 11.3, 0.7, "双 Harness 经验内核 —— 工程项目开发能力",
     size=24, color=INK2, bold=True)
text(s, 1.0, 4.35, 11.3, 0.5,
     "把 AI 编码 harness 从「无记忆的工具」，升级为带项目经验、带安全护栏、",
     size=15, color=INK3)
text(s, 1.0, 4.78, 11.3, 0.5,
     "经验可积累、可验证的工程开发系统。",
     size=15, color=INK3)
# 三个 hook 芯片
hook_x = 1.0
for label, c in [("pre_llm_call → recall", ACCENT),
                 ("on_session_end → session-end", VIOLET),
                 ("pre_tool_call → guard", RED)]:
    chip(s, hook_x, 5.55, 3.4, 0.5, label, color=c, size=11, fill=PANEL)
    hook_x += 3.62
# 底部 meta
text(s, 1.0, 6.55, 11.3, 0.4, "2026-08  ·  euly  ·  Hermes × Claude Code  ·  纯文件 + git 内核",
     size=12, color=INK3, font=MONO)

# ── S2 一句话定位 ──────────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "POSITIONING", "一句话定位", 2)
q = card(s, 0.75, 1.70, 11.83, 1.55, edge=ACC_D, fill=PANEL)
rect(s, 0.75, 1.70, 0.09, 1.55, fill=ACCENT)
tf = q.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.4)
p = para(tf, first=True, line=1.4)
run(p, "evo-kernel 是一个「", size=19, color=INK2)
run(p, "经验内核", size=19, color=INK, bold=True)
run(p, "」：纯文件 + git 存储的知识库，经 3 个 hook 挂到 AI 编码 harness，", size=19, color=INK2)
run(p, "把每一次对话自动转化为可检索、可复用、可固化的工程经验。", size=19, color=INK2)

stats = [("260+", "", "已验证条目\n(manifest)", ACCENT),
         ("6", " 类", "记忆形态\nfacts / playbook / skills…", VIOLET),
         ("3 × 2", "", "hooks × 双 harness\nHermes / Claude Code", GREEN)]
sx = 0.75
for num, unit, lab, c in stats:
    card(s, sx, 3.70, 3.77, 1.55)
    tf = tb(s, sx + 0.35, 3.90, 3.1, 0.72)
    p = para(tf, first=True)
    run(p, num, size=34, color=c, bold=True, font=MONO)
    if unit:
        run(p, unit, size=26, color=c, bold=True, font=CN)
    for i, ln in enumerate(lab.split("\n")):
        text(s, sx + 0.35, 4.62 + i * 0.28, 3.1, 0.3, ln, size=11.5, color=INK2)
    sx += 4.03
# 特征条
feat_y = 5.60
feats = [("经验即数据", "markdown + git 版本化，全程可审计"),
         ("零依赖内核", "单份 Node CLI（26 命令），无外部服务"),
         ("协议唯一", "只认 Claude hook 契约，换 harness 只换 adapter")]
fx = 0.75
for t, d in feats:
    card(s, fx, feat_y, 3.77, 1.15)
    text(s, fx + 0.3, feat_y + 0.18, 3.2, 0.3, t, size=14, color=INK, bold=True)
    text(s, fx + 0.3, feat_y + 0.55, 3.25, 0.55, d, size=11, color=INK3, line=1.15)
    fx += 4.03
footer(s, 2)

# ── S3 总体架构 ────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "ARCHITECTURE", "总体架构：一个内核，双 harness 接线", 3)
# 内核盒
kernel = card(s, 0.75, 1.62, 8.05, 2.35, edge=ACC_D, fill=PANEL)
text(s, 1.05, 1.78, 7.0, 0.35, "evo-kernel  经验内核", size=15, color=INK, bold=True, font=MONO)
text(s, 1.05, 2.14, 7.4, 0.3, "纯 markdown + git · Node CLI 26 命令 · 目录即状态机", size=10.5, color=INK3)
zones = [("inbox", "零判断暂存", INK3), ("lessons", "候选经验", INK3),
         ("facts", "语义记忆", VIOLET), ("episodes", "情景记忆", ACCENT),
         ("playbook", "策略库", GREEN), ("principles", "原则", CYAN),
         ("skills", "程序记忆", ACCENT), ("constraints", "硬约束", RED)]
zx, zy = 1.05, 2.60
for i, (name, role, c) in enumerate(zones):
    col = i % 4; row = i // 4
    xx = zx + col * 1.90; yy = zy + row * 0.62
    chip(s, xx, yy, 1.78, 0.46, name, color=c, size=10, fill=CARD)
# 三竖线 + 下箭头
for cx in (2.05, 4.77, 7.52):
    rect(s, cx, 3.99, 0.02, 0.32, fill=INK3)
    shape(s, MSO_SHAPE.DOWN_ARROW, cx - 0.10, 4.26, 0.22, 0.24, fill=INK3)
# adapter 层
adapt = card(s, 0.75, 4.56, 8.05, 1.28, edge=EDGE, fill=CARD)
text(s, 1.05, 4.70, 7.0, 0.3, "adapter 适配层   ~/.hermes/agent-hooks/", size=12.5, color=INK2, font=MONO)
text(s, 1.05, 4.99, 7.5, 0.3, "Hermes payload → evo 协议翻译 · 全部 fail-open（故障不阻塞主链路）", size=10, color=INK3)
chip(s, 1.05, 5.36, 2.35, 0.36, "evo-recall.sh", color=ACCENT, size=10)
chip(s, 3.55, 5.36, 2.45, 0.36, "evo-session-end.sh", color=VIOLET, size=10)
chip(s, 6.15, 5.36, 2.20, 0.36, "evo-guard.sh", color=RED, size=10)
# harness 层
text(s, 0.75, 6.12, 3.9, 0.5, "Hermes", size=16, color=INK, bold=True, font=MONO)
text(s, 0.75, 6.48, 3.9, 0.3, "config.yaml · hooks 段", size=10.5, color=INK3, font=MONO)
text(s, 4.90, 6.12, 3.9, 0.5, "Claude Code", size=16, color=INK, bold=True, font=MONO)
text(s, 4.90, 6.48, 3.9, 0.3, "settings.json · 同三件套", size=10.5, color=INK3, font=MONO)
rect(s, 0.75, 6.05, 8.05, 0.02, fill=EDGE)
# 右侧设计原则
px = 9.15
prin = [("协议唯一 · 适配器换装", "内核只认 Claude hook 契约；接入新 harness 只需写 3 个 adapter，双端共享同一份经验。", ACCENT),
        ("内核无状态", "零依赖 Node + markdown；git 即备份，也做审计。仓库位置无关。", VIOLET),
        ("双写通道", "会话结束双写：git 治理层（可审计）+ gbrain 语义脑（向量检索）。", GREEN)]
py = 1.62
for t, d, c in prin:
    card(s, px, py, 3.43, 1.70, edge=EDGE)
    rect(s, px, py, 0.07, 1.70, fill=c)
    text(s, px + 0.28, py + 0.16, 2.95, 0.35, t, size=13.5, color=INK, bold=True)
    text(s, px + 0.28, py + 0.58, 2.95, 1.0, d, size=10.5, color=INK2, line=1.25)
    py += 1.86
footer(s, 3)

# ── S4 三 hook 运行时 ──────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "RUNTIME", "三 hook：经验在对话中的实时闭环", 4)
hooks = [
    ("pre_llm_call", "→ recall", "检索注入", ACCENT,
     ["每次 LLM 调用前，按当前任务语义检索经验",
      "按 triggers / tags 计权，注入上下文",
      "让新对话「记得」所有历史教训"]),
    ("on_session_end", "→ session-end", "登记 + 双写", VIOLET,
     ["会话结束自动登记 transcript 引用",
      "双写 gbrain 语义脑（put_page + timeline）",
      "沉淀为可蒸馏的工程素材"]),
    ("pre_tool_call", "→ guard", "硬约束拦截", RED,
     ["工具调用前执行硬约束规则",
      "命中（如危险 rm -rf）即 deny 阻断",
      "经验从「建议」升级为「护栏」"]),
]
hx = 0.75
for ev, arrow, name, c, items in hooks:
    card(s, hx, 1.70, 3.77, 3.30, edge=EDGE)
    rect(s, hx, 1.70, 3.77, 0.09, fill=c)
    text(s, hx + 0.3, 1.95, 3.2, 0.35, ev + "  " + arrow, size=12.5, color=c, bold=True, font=MONO)
    text(s, hx + 0.3, 2.36, 3.2, 0.35, name, size=16, color=INK, bold=True)
    rect(s, hx + 0.3, 2.80, 3.17, 0.02, fill=EDGE)
    iy = 2.98
    for it in items:
        text(s, hx + 0.3, iy, 3.2, 0.6, "·  " + it, size=11.5, color=INK2, line=1.2)
        iy += 0.64
    hx += 4.03
note = card(s, 0.75, 5.35, 11.83, 0.95, edge=EDGE, fill=PANEL)
text(s, 1.05, 5.52, 11.2, 0.3, "I1 · fail-open", size=13, color=GREEN, bold=True, font=MONO)
text(s, 1.05, 5.82, 11.2, 0.4, "任何 hook / adapter 故障都不阻塞 harness 主链路 —— 经验系统是「增强」而非「依赖」，主链路永远可跑。", size=11.5, color=INK2)
footer(s, 4)

# ── S5 目录即状态机 ────────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "DATA MODEL", "目录即状态机：条目的位置 = 条目的状态", 5)
# 左：目录角色表
card(s, 0.75, 1.66, 6.6, 4.9, edge=EDGE, fill=PANEL)
text(s, 1.05, 1.84, 6.0, 0.3, "目录角色与注入资格（I2 守恒）", size=13.5, color=INK, bold=True)
rows = [("inbox/", "capture 暂存 + 会话登记", "永不注入", RED, False),
        ("lessons/", "candidate 经验暂存", "永不注入", RED, False),
        ("facts/", "语义记忆 · 事实/偏好/环境", "注入", GREEN, True),
        ("episodes/", "情景记忆 · 一次任务一份", "注入", GREEN, True),
        ("playbook/", "策略库 · 原子 bullet", "注入", GREEN, True),
        ("principles/", "原则 · 跨域普适", "注入", GREEN, True),
        ("skills/", "程序记忆 · SKILL.md 包", "skill 注入", GREEN, True),
        ("constraints/", "硬约束 · guard 执行", "实时拦截", RED, True),
        ("archive/", "退役/固化存档", "不检索", RED, False)]
ry = 2.24
for name, role, inj, c, active in rows:
    chip(s, 1.05, ry, 1.55, 0.34, name, color=ACCENT if active else INK3, size=10, fill=CARD)
    text(s, 2.72, ry, 2.7, 0.34, role, size=11, color=INK2)
    text(s, 5.55, ry, 1.5, 0.34, inj, size=11, color=c, bold=True, align=PP_ALIGN.RIGHT)
    ry += 0.465
# 右：状态迁移
card(s, 7.65, 1.66, 4.93, 4.9, edge=EDGE, fill=PANEL)
text(s, 7.95, 1.84, 4.3, 0.3, "状态迁移", size=13.5, color=INK, bold=True)
flow = [("inbox", "→ curate →", "facts / episodes / playbook / principles", "人审唯一入库口"),
        ("validated", "→ solidify →", "skills（程序记忆）/ constraints（硬约束）", "反复验证后固化"),
        ("validated", "→ demote →", "archive（退役）/ lessons（回炉）", "失效条目退出注入")]
fy = 2.22
for a, b, c_, d in flow:
    tf = tb(s, 7.95, fy, 4.3, 0.34)
    p = para(tf, first=True)
    run(p, a + "  ", size=11.5, color=INK3, font=MONO)
    run(p, b, size=12, color=ACCENT, bold=True, font=MONO)
    text(s, 7.95, fy + 0.36, 4.3, 0.32, c_, size=11.5, color=INK)
    text(s, 7.95, fy + 0.68, 4.3, 0.28, d, size=10, color=INK3)
    fy += 1.18
rect(s, 7.95, fy + 0.04, 4.3, 0.02, fill=EDGE)
text(s, 7.95, fy + 0.18, 4.3, 0.3, "I3 · 人审前置  —  curate 是唯一入库口", size=10.5, color=AMBER, bold=True)
footer(s, 5)

# ── S6 经验生命周期 ────────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "LIFECYCLE", "一条经验的完整旅程", 6)
stages = [("捕获", "零判断入口"),
          ("登记", "自动登记会话"),
          ("蒸馏", "切硬证据"),
          ("人审", "唯一入库口"),
          ("注入", "动手前在场"),
          ("验证", "对账回填"),
          ("固化/退役", "skill / hook")]
sx = 0.55; sw = 1.62; gap = 0.16
for i, (t, d1) in enumerate(stages):
    x = sx + i * (sw + gap)
    shape(s, MSO_SHAPE.CHEVRON, x, 1.80, sw, 1.55, fill=CHEV, line=EDGE, lw=1.0)
    text(s, x + 0.14, 2.14, sw - 0.40, 0.4, t, size=13.5, color=INK, bold=True)
    text(s, x + 0.14, 2.62, sw - 0.40, 0.5, d1, size=9.5, color=INK2, line=1.1)
# 治理不变量条
inv = [("I1 · fail-open", "故障不阻塞主链路", GREEN),
       ("I2 · 注入守恒", "未审条目永不注入", ACCENT),
       ("I3 · 人审前置", "curate 是唯一入库口", AMBER),
       ("I4 · 计数对账", "helpful/harmful 单点回填", VIOLET)]
ix = 0.75
for t, d, c in inv:
    card(s, ix, 3.85, 2.85, 1.05, edge=EDGE)
    text(s, ix + 0.25, 4.00, 2.45, 0.3, t, size=12.5, color=c, bold=True)
    text(s, ix + 0.25, 4.32, 2.45, 0.5, d, size=10, color=INK3, line=1.1)
    ix += 3.0
# 底部闭环示意
text(s, 0.75, 5.45, 11.8, 0.3, "为什么有效：条目不随会话蒸发", size=13, color=INK, bold=True)
text(s, 0.75, 5.80, 11.8, 0.9,
     "经验以「候选 → 验证 → 固化」递进；验证靠实战计数（helpful/harmful），差经验被审计与退役，库里长期留存的都是经实战验证的结论。",
     size=12, color=INK2, line=1.3)
footer(s, 6)

# ── S7 六类记忆（工程开发能力） ────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "CAPABILITY", "六类记忆：工程能力的沉淀载体", 7)
mem = [("episodes", "情景记忆", "6", "一次任务一份工程档案", "任务级结论 / 决策 / 证据", ACCENT),
       ("facts", "语义记忆", "18", "项目环境与已验证结论", "按 domains 分域沉淀", VIOLET),
       ("playbook", "策略库", "90", "原子操作手册", "每条带 3-5 个 triggers", GREEN),
       ("principles", "原则", "3", "跨域普适开发纪律", "方法论层面，适用所有项目", CYAN),
       ("skills", "程序记忆", "5", "SKILL.md 可直接执行", "反复验证后固化", AMBER),
       ("constraints", "硬约束", "2", "guard 实时拦截", "危险操作机器强制执行", RED)]
mx, my = 0.75, 1.70
for i, (en, cn, num, d1, d2, c) in enumerate(mem):
    col = i % 3; row = i // 3
    x = mx + col * 4.03; y = my + row * 2.55
    card(s, x, y, 3.77, 2.30, edge=EDGE)
    rect(s, x, y, 3.77, 0.08, fill=c)
    text(s, x + 0.3, y + 0.22, 2.4, 0.4, en, size=14, color=c, bold=True, font=MONO)
    text(s, x + 2.15, y + 0.20, 1.35, 0.5, num, size=30, color=INK, bold=True, font=MONO, align=PP_ALIGN.RIGHT)
    text(s, x + 0.3, y + 0.72, 3.2, 0.35, cn, size=14, color=INK, bold=True)
    text(s, x + 0.3, y + 1.08, 3.2, 0.35, d1, size=11.5, color=INK2)
    text(s, x + 0.3, y + 1.44, 3.2, 0.7, d2, size=10, color=INK3, line=1.2)
footer(s, 7)

# ── S8 案例特写：UCM221 雷达 ───────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "CASE STUDY", "案例特写：UCM221 雷达项目", 8)
case = card(s, 0.75, 1.70, 6.9, 3.9, edge=ACC_D, fill=PANEL)
text(s, 1.05, 1.88, 6.3, 0.3, "episode · faf-legacy-gate-domain-bugs", size=12, color=ACCENT, font=MONO)
text(s, 1.05, 2.22, 6.3, 0.75, "FAF 移植发现两个静默门限 bug", size=20, color=INK, bold=True)
text(s, 1.05, 3.05, 6.3, 1.0,
     "俯仰门限作用在错误的量上（弧度 vs sin），方位门限是从未生效的死代码。",
     size=13.5, color=INK2, line=1.35)
rect(s, 1.05, 4.05, 6.3, 0.02, fill=EDGE)
text(s, 1.05, 4.20, 6.3, 0.3, "效果可量化", size=12, color=GREEN, bold=True)
text(s, 1.05, 4.52, 6.3, 0.5, "航迹覆盖 99.7% → 100.0%  ·  误差 0.18 m → 0.12 m", size=14, color=INK, bold=True, font=MONO)
text(s, 1.05, 5.06, 6.3, 0.4, "下次任何人再碰 FAF 门限，这些坑在动手前已注入。", size=12, color=INK3)
# 右：playbook 弹药库
card(s, 7.95, 1.70, 4.63, 3.9, edge=EDGE, fill=PANEL)
text(s, 8.25, 1.88, 4.0, 0.3, "playbook · 策略弹药库", size=13, color=INK, bold=True)
pb = [("chamber-doa-bin3-fixed-gate-rx-pair", "暗室测角固定门限的 RX 配对方案"),
      ("radar-doppler-bin-wraparound-unroll-first", "多普勒 bin 卷绕先展开再处理"),
      ("armv7-cross-compile-needs-mfpu-neon", "ARMv7 交叉编译必须显式开 NEON")]
py = 2.30
for name, d in pb:
    chip(s, 8.25, py, 4.0, 0.52, name, color=GREEN, size=9.5, fill=CARD)
    text(s, 8.30, py + 0.55, 4.0, 0.45, d, size=10, color=INK3, line=1.1)
    py += 1.12
# 底部事实条
card(s, 0.75, 5.85, 11.83, 1.0, edge=EDGE, fill=CARD)
text(s, 1.05, 6.00, 11.2, 0.3, "facts · 项目现状随手可查", size=12.5, color=VIOLET, bold=True, font=MONO)
text(s, 1.05, 6.30, 11.2, 0.4,
     "交叉角数据集档案 · 仓结构重构记录 · tracker 宏定义位置 · ARM 超分辨性能结论 —— 新会话进来项目「现状」已在场。",
     size=11.5, color=INK2)
footer(s, 8)

# ── S9 闭环与三大能力 ──────────────────────────────────────────────────────
s = new_slide(); bg(s)
header(s, "CLOSED LOOP", "闭环效果：三大工程能力", 9)
loop = ["会话发生", "自动登记", "周期蒸馏", "人审入库", "自动注入", "对账验证", "固化/退役"]
lx = 0.55; lw = 1.62; lg = 0.16
for i, t in enumerate(loop):
    x = lx + i * (lw + lg)
    shape(s, MSO_SHAPE.CHEVRON, x, 1.80, lw, 0.95, fill=CARD, line=EDGE, lw=1.0)
    text(s, x + 0.12, 2.10, lw - 0.2, 0.4, t, size=12, color=INK, bold=True, align=PP_ALIGN.CENTER)
ret = card(s, 0.75, 3.05, 11.83, 0.62, edge=EDGE, fill=PANEL)
text(s, 1.0, 3.20, 11.4, 0.35, "⟲  下一次会话开始时，历史经验已在场 —— 回到「会话发生」", size=12.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
caps = [("项目知识不随会话蒸发", "每次对话结论落为持久化条目；双 harness 共享同一份库，跨工具、跨天、跨会话连续。", ACCENT),
        ("经验在最需要的时刻出现", "recall 按 triggers 语义匹配注入；不必主动检索，相关教训动手前已在上下文。", GREEN),
        ("质量有账可查", "每条经验带 helpful/harmful 计数与验证等级；差经验被统计、审计、退役。", VIOLET)]
cx = 0.75
for t, d, c in caps:
    card(s, cx, 4.05, 3.77, 2.35, edge=EDGE)
    rect(s, cx, 4.05, 3.77, 0.08, fill=c)
    text(s, cx + 0.3, 4.32, 3.2, 0.6, t, size=15, color=INK, bold=True, line=1.15)
    text(s, cx + 0.3, 5.00, 3.2, 1.2, d, size=11, color=INK2, line=1.3)
    cx += 4.03
footer(s, 9)

# ── S10 总结 ───────────────────────────────────────────────────────────────
s = new_slide(); bg(s)
rect(s, 0, 7.34, 13.333, 0.16, fill=ACC_D)
text(s, 1.0, 1.7, 11.3, 0.35, "SUMMARY", size=13, color=ACCENT, font=MONO)
text(s, 0.95, 2.15, 11.4, 1.1, "经验内核，让 AI 开发越用越强", size=40, color=INK, bold=True)
rect(s, 1.0, 3.35, 0.7, 0.06, fill=GREEN)
text(s, 1.0, 3.62, 11.3, 0.45,
     "把 AI 编码 harness 从「无记忆的工具」，升级为带项目经验、带安全护栏、",
     size=17, color=INK2)
text(s, 1.0, 4.14, 11.3, 0.45,
     "经验可积累、可验证的工程开发系统。",
     size=17, color=INK2)
vals = ["经验即数据 · git 可审计", "三 hook 实时闭环", "人审前置 + fail-open", "六类记忆 + 硬约束护栏"]
vx = 1.0
for v in vals:
    chip(s, vx, 5.15, 2.75, 0.5, v, color=ACCENT, size=11, fill=PANEL)
    vx += 2.88
text(s, 1.0, 6.35, 11.3, 0.4,
     "evo-kernel  ·  ~/Dev/evo-kernel  ·  纯文件 + git  ·  Hermes × Claude Code",
     size=12, color=INK3, font=MONO)

OUT = sys.argv[1] if len(sys.argv) > 1 else "/Users/zodyne/Dev/evo-kernel/docs/evo-kernel-architecture.pptx"
prs.save(OUT)
print("saved:", OUT, "| slides:", len(prs.slides._sldIdLst))
